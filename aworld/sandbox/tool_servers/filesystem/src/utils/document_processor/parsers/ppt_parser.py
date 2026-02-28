"""PowerPoint 解析器（需 LibreOffice + pdf2image 生成截图）"""

import base64
import logging
import shutil
import subprocess
import sys
from pathlib import Path
from typing import Any, Dict, List, Optional

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from ..base_parser import BaseParser

logger = logging.getLogger(__name__)


class PptParser(BaseParser):
    def get_supported_types(self) -> list[str]:
        return ["pptx", "ppt"]

    def can_handle(self, file_type: str) -> bool:
        return file_type.lower() in ["pptx", "ppt"]

    async def parse(
        self,
        file_path: Path,
        task_id: str = "",
        source_file_id: Optional[str] = None,
        source_file_name: Optional[str] = None,
        afts_service: Any = None,
        output_path: Optional[Path] = None,
        **kwargs,
    ) -> Dict[str, Any]:
        if Presentation is None:
            raise RuntimeError("未安装 python-pptx。请安装: pip install python-pptx")
        out_path = Path(output_path) if output_path else None
        images_dir = (out_path.parent / "images") if out_path else (
            Path.home() / "aworld_workspace" / (task_id or "out") / "images"
        )
        screenshot_paths = await self._generate_slide_screenshots(file_path, images_dir)
        embed = kwargs.get("embed_images", False)
        markdown_content = await self._extract_pptx_content(file_path, images_dir, screenshot_paths, embed)
        if out_path:
            out_path.parent.mkdir(parents=True, exist_ok=True)
            out_path.write_text(markdown_content, encoding="utf-8")
            return {"file_path": out_path}
        source_name = source_file_name or file_path.stem
        parsed_file_path = await self._save_markdown_to_file(
            markdown_content, task_id, source_name, output_path=output_path
        )
        return {"file_path": parsed_file_path}

    def _find_libreoffice(self) -> Optional[str]:
        if sys.platform == "win32":
            candidates = ["soffice"]
        elif sys.platform == "darwin":
            candidates = [
                "/Applications/LibreOffice.app/Contents/MacOS/soffice",
                "/usr/local/bin/libreoffice",
                "/opt/homebrew/bin/libreoffice",
                "soffice",
            ]
        else:
            candidates = ["libreoffice", "soffice"]
        for cmd in candidates:
            if shutil.which(cmd) or Path(cmd).exists():
                try:
                    subprocess.run([cmd, "--version"], capture_output=True, check=True, timeout=5)
                    return cmd
                except (subprocess.CalledProcessError, FileNotFoundError, subprocess.TimeoutExpired):
                    continue
        return None

    async def _generate_slide_screenshots(self, file_path: Path, images_dir: Path) -> List[Path]:
        libreoffice_cmd = self._find_libreoffice()
        if not libreoffice_cmd:
            logger.warning("未找到 LibreOffice，跳过截图")
            return []
        images_dir.mkdir(parents=True, exist_ok=True)
        presentation = Presentation(str(file_path))
        num_slides = len(presentation.slides)
        screenshot_paths = []
        temp_pdf = images_dir / f"{file_path.stem}_temp.pdf"
        try:
            subprocess.run(
                [libreoffice_cmd, "--headless", "--convert-to", "pdf", "--outdir", str(images_dir), str(file_path)],
                capture_output=True,
                text=True,
                timeout=120,
            )
            pdf_file = images_dir / f"{file_path.stem}.pdf"
            if not pdf_file.exists():
                return []
            pdf_file.rename(temp_pdf)
            try:
                from pdf2image import convert_from_path
            except ImportError:
                if temp_pdf.exists():
                    temp_pdf.unlink(missing_ok=True)
                return []
            images = convert_from_path(str(temp_pdf), dpi=150, first_page=1, last_page=num_slides)
            for i, img in enumerate(images, 1):
                p = images_dir / f"slide_{i:02d}.png"
                img.save(p, "PNG")
                screenshot_paths.append(p)
        except Exception as e:
            logger.warning("ppt_parser _generate_slide_screenshots: %s", e)
        finally:
            if temp_pdf.exists():
                temp_pdf.unlink(missing_ok=True)
        return sorted(screenshot_paths, key=lambda p: p.name)

    async def _extract_pptx_content(
        self,
        file_path: Path,
        images_dir: Path,
        screenshot_paths: Optional[List[Path]] = None,
        embed_images: bool = False,
    ) -> str:
        presentation = Presentation(str(file_path))
        parts = []
        for slide_idx, slide in enumerate(presentation.slides):
            n = slide_idx + 1
            parts.append(f"## 幻灯片 {n}\n\n")
            expected_name = f"slide_{n:02d}.png"
            screenshot_path = images_dir / expected_name
            rel_path = Path("images") / expected_name
            if screenshot_path.exists():
                if embed_images:
                    try:
                        b64 = base64.b64encode(screenshot_path.read_bytes()).decode("ascii")
                        parts.append(f"![幻灯片 {n} 截图](data:image/png;base64,{b64})\n\n")
                    except Exception:
                        parts.append(f"![幻灯片 {n} 截图]({rel_path})\n\n")
                else:
                    parts.append(f"![幻灯片 {n} 截图]({rel_path})\n\n")
            else:
                parts.append(f"![幻灯片 {n} 截图]({rel_path})\n\n")
            parts.append("### 文本内容\n\n")
            has_text = False
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    has_text = True
                    text = shape.text.strip()
                    is_title = False
                    try:
                        if getattr(shape, "is_placeholder", False) and getattr(shape, "placeholder_format", None):
                            if "title" in str(shape.placeholder_format.type).lower():
                                is_title = True
                    except Exception:
                        pass
                    if is_title:
                        parts.append(f"**标题**: {text}\n\n")
                    else:
                        parts.append(f"{text}\n\n")
                if hasattr(shape, "table"):
                    try:
                        t = shape.table
                        if t.rows:
                            parts.append("\n**表格**:\n\n")
                            headers = [c.text.strip() for c in t.rows[0].cells if c.text.strip()]
                            if headers:
                                parts.append("| " + " | ".join(headers) + " |\n")
                                parts.append("| " + " | ".join(["---"] * len(headers)) + " |\n")
                            for row in t.rows[1:]:
                                parts.append(
                                    "| "
                                    + " | ".join(c.text.strip() for c in row.cells if c.text.strip())
                                    + " |\n"
                                )
                            parts.append("\n")
                    except Exception:
                        pass
            if not has_text:
                parts.append("（无文本内容）\n\n")
            if getattr(slide, "notes_slide", None) and getattr(slide.notes_slide, "notes_text_frame", None):
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    parts.append(f"**备注**: {notes}\n\n")
            parts.append("---\n\n")
        return "".join(parts)

